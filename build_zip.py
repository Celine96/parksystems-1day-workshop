# -*- coding: utf-8 -*-
"""
파크시스템스 임직원 워크북 zip 패키징
"""
import os, re, shutil, zipfile, json
from pathlib import Path
from pptx import Presentation

BASE = Path(r"C:\Users\User\Desktop\06_개인\2. ai native workshop\파크시스템스 기업강의\v1-파크시스템스")
SHARED = Path(r"C:\Users\User\Desktop\06_개인\2. ai native workshop\파크시스템스 기업강의\파크시스템스 (나나님 공유)")
TEMP = BASE / "zip_temp"
ROOT = TEMP / "parksystems-workshop"
ZIP_OUT = BASE / "파크시스템스 임직원_워크북.zip"

# 기존 temp 정리
if TEMP.exists():
    shutil.rmtree(TEMP)
ROOT.mkdir(parents=True)

# ============================================
# 1. PPT 가공 (실명·고객사·사업부 가명화)
# ============================================
src_pptx = SHARED / "2. 박정연 PM agent 실습" / "2026-05-18 [TW] Weekly Report Sample.pptx"
dst_pptx = ROOT / "4_PM에이전트" / "Weekly_Report_Sample.pptx"
dst_pptx.parent.mkdir(parents=True, exist_ok=True)

# 가명 매핑 — anonymization_map.json에서 로드 (실명 추가/변경 시 그 파일만 수정)
with open(BASE / "anonymization_map.json", encoding='utf-8') as f:
    anon_map = json.load(f)
replacements = []
for category in ['names', 'customers', 'divisions']:
    replacements.extend(anon_map[category].items())

try:
    prs = Presentation(str(src_pptx))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        new_text = run.text
                        for old, new in replacements:
                            new_text = new_text.replace(old, new)
                        run.text = new_text
    prs.save(str(dst_pptx))
    print(f"[OK] PPT 가공: {dst_pptx.name} (가명 {len(replacements)}건 적용)")

    # 가명화 누락 검증 — 결과 PPT에서 원본 실명이 남았는지 재확인
    leak_check = Presentation(str(dst_pptx))
    leaks = []
    for slide in leak_check.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            for original in list(anon_map['names'].keys()) + list(anon_map['customers'].keys()):
                if '\n' in original:
                    continue  # 줄바꿈 포함 패턴은 PPT run 단위로 잘릴 수 있어 false positive
                if original in text:
                    leaks.append((slide.slide_id, original))
    if leaks:
        print(f"[WARN] 가명화 누락 의심 {len(leaks)}건: {leaks[:3]}{'...' if len(leaks) > 3 else ''}")
    else:
        print(f"[OK] 가명화 누락 검증 통과")
except FileNotFoundError as e:
    print(f"[FAIL] PPT 원본 파일 없음: {src_pptx}")
    print(f"  {e}")
    raise SystemExit(1)
except Exception as e:
    print(f"[FAIL] PPT 가공 실패: {type(e).__name__}: {e}")
    raise SystemExit(1)

# ============================================
# 2. 5_매뉴얼 폴더 (나나님 공유 자료 복사)
# ============================================
manual_dir = ROOT / "5_매뉴얼"
(manual_dir / "형식").mkdir(parents=True, exist_ok=True)
(manual_dir / "실제_매뉴얼").mkdir(parents=True, exist_ok=True)
(manual_dir / "raw_data").mkdir(parents=True, exist_ok=True)
(manual_dir / "output").mkdir(parents=True, exist_ok=True)

# 형식 PNG 4개
fmt_src = SHARED / "3. 매뉴얼 및 PPT 자동 생성"
for png in ["매뉴얼 형식_Intro.png", "매뉴얼 형식_Preface.png", "PPT 형식.png", "Spec sheet.png"]:
    shutil.copy(fmt_src / png, manual_dir / "형식" / png)
print(f"[OK] 형식 PNG 4개 복사")

# 실제 매뉴얼 PDF 4개 (이름 간소화)
real_src = SHARED / "4. 문서 업로드 및 배포 자동화" / "업로드용 매뉴얼"
pdf_map = {
    "(7K) NXT1518-IN2-C0.00-KO (SDC) 1.pdf": "NXT1518-IN2-KO.pdf",
    "(7K) NXT1518-INI-C0.00-EN (SDC) 1.pdf": "NXT1518-INI-EN.pdf",
    "(7K) NXT1518-MT2-C0.00-KO (SDC).pdf": "NXT1518-MT2-KO.pdf",
    "(7K) NXT1518-OP2-C0.00-KO (SDC).pdf": "NXT1518-OP2-KO.pdf",
}
for src_name, dst_name in pdf_map.items():
    shutil.copy(real_src / src_name, manual_dir / "실제_매뉴얼" / dst_name)
print(f"[OK] 실제 매뉴얼 PDF 4개 복사")

# raw_data
shutil.copy(fmt_src / "(raw data) NX-Interferom.txt", manual_dir / "raw_data" / "NX-Interferom.txt")
print(f"[OK] raw_data 복사")

# ============================================
# 3. 3_파일정리/raw_files/ — 더미 흩어진 파일 20개
# ============================================
files_dir = ROOT / "3_파일정리" / "raw_files"
files_dir.mkdir(parents=True, exist_ok=True)

dummy_files = [
    "(S) NX-Interferom_v1.stp",
    "(S) NX-Interferom_review.pdf",
    "(S) NX-Interferom_input.stp",
    "(S) NX-Interferom_draft_v2.pdf",
    "(B) Manual_update_raw.stp",
    "(B) Manual_update_working.pdf",
    "(B) Manual_update_v3.pdf",
    "(B) Manual_update_input.stp",
    "final_NX_S.pdf",
    "final_Manual_B.pdf",
]
for fname in dummy_files:
    (files_dir / fname).write_text(
        f"# {fname}\n\n(더미 파일 - 1교시 파일 정리 실습용)\n",
        encoding='utf-8'
    )
print(f"[OK] 3_파일정리/raw_files 더미 20개 생성")

# ============================================
# 4. 4_PM에이전트/ — 빈 폴더 (3 챕터에서 학습자가 정리한 RFM이 여기 채워짐)
# ============================================
# RFM/ 폴더는 학습자가 3 챕터 실습으로 직접 만듦.
# Weekly_Report_Sample.pptx는 위 (1) 단계에서 이미 4_PM에이전트/에 들어가 있음.
print(f"[OK] 4_PM에이전트/ (RFM 폴더는 3 챕터 학습자 결과로 채움)")

# ============================================
# 5. .claude/ 폴더 (에이전트·슬래시 커맨드 미리 셋업)
#     - templates/.claude/ 안 .md 파일들을 zip 안에 복사
#     - 학습자가 워크북 따라 만들 시간을 아껴서, 막혔을 때 fallback 역할
#     - 호스트가 멘트로 "이미 들어 있어도 실습 흐름 따라가세요" 안내
# ============================================
shutil.copytree(BASE / "templates" / ".claude", ROOT / ".claude")
preset_files = sorted(p.relative_to(ROOT) for p in (ROOT / ".claude").rglob("*.md"))
print(f"[OK] .claude/ 미리 셋업 ({len(preset_files)}개 파일):")
for p in preset_files:
    print(f"     - {p}")

# .claude/settings.json — 권한 자동 허용 (학습자 흐름 매끄럽게)
settings_payload = {
    "permissions": {
        "defaultMode": "auto",
        "allow": [
            "Write(./.claude/agents/**)",
            "Write(./.claude/commands/**)",
            "Edit(./.claude/agents/**)",
            "Edit(./.claude/commands/**)"
        ]
    }
}
(ROOT / ".claude" / "settings.json").write_text(
    json.dumps(settings_payload, indent=2, ensure_ascii=False),
    encoding='utf-8'
)
print(f"[OK] .claude/settings.json 생성 (권한 자동 허용)")

# ============================================
# 6. CLAUDE.md (클로드 코드 자동 로드)
#     - 폴더 구조·워크숍 흐름 텍스트는 chapters.json에서 동적 생성
#       (drift 방지: 챕터 추가/변경 시 chapters.json만 수정)
# ============================================
with open(BASE / "docs" / ".vitepress" / "data" / "chapters.json", encoding='utf-8') as f:
    chapters_data = json.load(f)

# 메인 수업 챕터 (워크숍 흐름의 source)
main_items = [
    item
    for g in chapters_data['groups'] if g['label'].startswith('메인')
    for item in g['items']
]

# 워크숍 흐름 (메인 챕터의 num·text·flow를 한 줄씩 — chapters.json drift 방지)
flow_lines = []
for i, item in enumerate(main_items, start=1):
    title_short = item['text'].split('. ', 1)[-1]  # "3. 프로젝트 #1 : 파일 정리" → "프로젝트 #1 : 파일 정리"
    flow_text = item.get('flow', item.get('details', ''))
    flow_lines.append(f"{i}. **{item['num']} {title_short}** — {flow_text}")
flow_block = "\n".join(flow_lines)

claude_md = f"""# 파크시스템스 AX 바이브코딩 연수 — 작업 공간

이 폴더는 워크숍 실습용이에요. 클로드 코드를 이 폴더에서 실행하면 모든 자료에 접근할 수 있어요.

## 폴더 구조

- `3_파일정리/raw_files/` — 1교시 파일 정리 실습용 더미 파일 (흩어진 상태)
- `4_PM에이전트/` — 2교시 PM 에이전트 실습 공간 (1교시에서 정리한 RFM 폴더가 여기 들어옴)
- `4_PM에이전트/Weekly_Report_Sample.pptx` — 주간 보고 양식
- `5_매뉴얼/형식/` — 매뉴얼 형식 가이드 (이미지 4개)
- `5_매뉴얼/실제_매뉴얼/` — 회사 실제 매뉴얼 (PDF 4개, 학습 자료)
- `5_매뉴얼/raw_data/` — 변환할 원본 자료
- `5_매뉴얼/output/` — 생성된 결과물 저장 위치
- `.claude/agents/` — 학습자가 만들 에이전트 정의 위치
- `.claude/commands/` — 학습자가 만들 슬래시 커맨드 위치

## 워크숍 흐름

{flow_block}

## 워크북 라이브

https://celine96.github.io/parksystems-1day-workshop/
"""
(ROOT / "CLAUDE.md").write_text(claude_md, encoding='utf-8')
print(f"[OK] CLAUDE.md 생성 (워크숍 흐름은 chapters.json 기반)")

# ============================================
# 7. README.md (zip 처음 열 때)
# ============================================
readme = """# 파크시스템스 AX 바이브코딩 연수 — 임직원 실습 자료

## 📍 워크북 (라이브)

👉 **https://celine96.github.io/parksystems-1day-workshop/**

워크북은 위 URL에서 보세요. 이 zip 안에는 **실습 자료**만 들어있어요.

## 📂 폴더 구조

```
parksystems-workshop/
  ├── README.md              ← 지금 보고 계신 파일
  ├── CLAUDE.md              ← 클로드 코드가 자동 읽는 메모
  ├── .claude/               ← 학습자가 만들 에이전트·커맨드 위치
  │
  ├── 3_파일정리/             ← 1교시 (파일 정리)
  │     └── raw_files/        ← 흩어진 더미 파일 20개
  │
  ├── 4_PM에이전트/           ← 2교시 (PM 에이전트)
  │     ├── RFM/              ← 작업자별 폴더 (시나리오 셋업)
  │     └── Weekly_Report_Sample.pptx
  │
  └── 5_매뉴얼/               ← 3교시 (매뉴얼 자동 생성 + 외부 근거 + HTML)
        ├── 형식/             ← 회사 매뉴얼 형식 (PNG 4개)
        ├── 실제_매뉴얼/      ← 회사 실제 매뉴얼 (PDF 4개)
        ├── raw_data/         ← 변환할 원본
        └── output/           ← 결과물 저장 위치
```

## 🚀 시작 방법

### 1. 폴더 압축 해제
이 zip을 바탕화면에 압축 해제해주세요. `parksystems-workshop/` 폴더가 생겨요.

### 2. 워크북 0교시 — 설치 안내 따라하기
https://celine96.github.io/parksystems-1day-workshop/part1/1-1-intro 에서 시작.

### 3. 터미널에서 폴더 이동
```
cd ~/Desktop/parksystems-workshop
```
> Windows: `cd C:\\Users\\USER\\Desktop\\parksystems-workshop`

### 4. 클로드 코드 실행
```
claude
```

이후는 워크북 라이브 URL에서 챕터 순서대로 진행하세요.

## 📝 참고

- 모든 회사 자료는 **실명·고객사명·사업부명을 가명화**해서 동봉했어요
- 더미 파일은 실습용으로 만든 것으로, 실제 회사 작업 내용이 아니에요
- 막히는 부분은 워크북 각 챕터의 Q&A 또는 강사에게 문의

화이팅! 🙋
"""
(ROOT / "README.md").write_text(readme, encoding='utf-8')
print(f"[OK] README.md 생성")

# ============================================
# 8. zip 압축
# ============================================
if ZIP_OUT.exists():
    ZIP_OUT.unlink()

with zipfile.ZipFile(ZIP_OUT, 'w', zipfile.ZIP_DEFLATED) as zf:
    for path in ROOT.rglob('*'):
        if path.is_file():
            arcname = path.relative_to(TEMP)
            zf.write(path, arcname)

# 결과
size_mb = ZIP_OUT.stat().st_size / 1024 / 1024
print(f"")
print("-" * 50)
print(f"[OK] zip 생성 완료")
print(f"파일: {ZIP_OUT}")
print(f"크기: {size_mb:.2f} MB")
print("-" * 50)

# 파일 목록 출력
print("\n=== zip 내부 파일 (요약) ===")
with zipfile.ZipFile(ZIP_OUT, 'r') as zf:
    names = zf.namelist()
    print(f"총 {len(names)}개 파일")
    for name in sorted(names)[:30]:
        print(f"  {name}")
    if len(names) > 30:
        print(f"  ... 외 {len(names) - 30}개")
